#!/usr/bin/env python3
"""
Convert PPTX/PPT to PNG slide images using python-pptx + pdf2image
Usage: python3 pptx_to_png.py <input.pptx> <output_dir>
Prints JSON: {"success": true, "slides": [...], "count": N}
"""
import sys, os, json, subprocess, tempfile, shutil

def main():
    if len(sys.argv) < 3:
        print(json.dumps({"success": False, "error": "Usage: pptx_to_png.py <input> <outdir>"}))
        sys.exit(1)

    input_path = sys.argv[1]
    out_dir    = sys.argv[2]
    os.makedirs(out_dir, exist_ok=True)

    ext = os.path.splitext(input_path)[1].lower()

    # Step 1: PPTX → PDF using python-pptx slide dimensions + reportlab, or
    # better: use LibreOffice if available, else python-pptx direct rendering
    pdf_path = None

    if ext in ('.pptx', '.ppt', '.odp'):
        # Try python-pptx + subprocess to convert to PDF
        # Method: use comtypes on Windows, or unoconv/libreoffice on Linux
        # Best cross-platform: use python-pptx to extract slide info, then
        # render each slide to an image using Pillow

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from PIL import Image, ImageDraw, ImageFont
            import io

            prs = Presentation(input_path)
            slide_w = int(prs.slide_width)
            slide_h = int(prs.slide_height)
            # Target render size
            scale = 1920 / slide_w
            W = 1920
            H = int(slide_h * scale)

            slides_out = []
            for i, slide in enumerate(prs.slides):
                img = Image.new('RGB', (W, H), color=(10, 10, 30))
                draw = ImageDraw.Draw(img)

                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    # Get position/size
                    left  = int(shape.left  * scale) if shape.left  else 0
                    top   = int(shape.top   * scale) if shape.top   else 0
                    width = int(shape.width * scale) if shape.width else W
                    height= int(shape.height* scale) if shape.height else H//4

                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if not text:
                            continue
                        # Try to get font size
                        fs = 36
                        color = (255, 255, 255)
                        for run in para.runs:
                            if run.font.size:
                                fs = max(12, int(run.font.size.pt * scale * 0.8))
                            if run.font.color and run.font.color.type:
                                try:
                                    rgb = run.font.color.rgb
                                    color = (rgb.r, rgb.g, rgb.b)
                                except:
                                    pass
                        try:
                            font = ImageFont.truetype('/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf', fs)
                        except:
                            font = ImageFont.load_default()
                        draw.text((left + 10, top), text, fill=color, font=font)
                        top += fs + 4

                out_path = os.path.join(out_dir, f'slide-{i+1}.png')
                img.save(out_path, 'PNG')
                slides_out.append({'index': i, 'imagePath': out_path})

            print(json.dumps({"success": True, "slides": slides_out, "count": len(slides_out)}))
            return

        except ImportError as e:
            # PIL not available, fall through to LibreOffice
            pass

    # Fallback: try LibreOffice
    lo_paths = [
        '/usr/bin/libreoffice', '/usr/bin/soffice',
        r'C:\Program Files\LibreOffice\program\soffice.exe',
        r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',
    ]
    lo = next((p for p in lo_paths if os.path.exists(p)), None)

    if lo and ext in ('.pptx', '.ppt', '.odp'):
        tmp = tempfile.mkdtemp()
        try:
            subprocess.run([lo, '--headless', '--convert-to', 'pdf', '--outdir', tmp, input_path],
                           timeout=120, check=True, capture_output=True)
            pdfs = [f for f in os.listdir(tmp) if f.endswith('.pdf')]
            if pdfs:
                pdf_path = os.path.join(tmp, pdfs[0])
        except Exception as e:
            shutil.rmtree(tmp, ignore_errors=True)
            print(json.dumps({"success": False, "error": f"Conversion failed: {e}"}))
            return
    elif ext == '.pdf':
        pdf_path = input_path

    if not pdf_path or not os.path.exists(pdf_path):
        print(json.dumps({"success": False, "error": "Could not produce PDF from input"}))
        return

    # Rasterize PDF with pdftoppm
    slide_prefix = os.path.join(out_dir, 'slide')
    try:
        subprocess.run(['pdftoppm', '-png', '-r', '150', pdf_path, slide_prefix],
                       timeout=120, check=True, capture_output=True)
    except Exception as e:
        print(json.dumps({"success": False, "error": f"pdftoppm failed: {e}"}))
        return

    slides_out = []
    files = sorted(
        [f for f in os.listdir(out_dir) if f.startswith('slide') and f.endswith('.png')],
        key=lambda f: int(''.join(filter(str.isdigit, f)) or '0')
    )
    for i, f in enumerate(files):
        slides_out.append({'index': i, 'imagePath': os.path.join(out_dir, f)})

    if not slides_out:
        print(json.dumps({"success": False, "error": "No slides generated"}))
        return

    print(json.dumps({"success": True, "slides": slides_out, "count": len(slides_out)}))

if __name__ == '__main__':
    main()
